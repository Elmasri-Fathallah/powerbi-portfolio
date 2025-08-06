from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Define colors (Saudi colors with Power BI accent)
green_primary = RGBColor(27, 94, 32)  # Dark green
gold_accent = RGBColor(255, 215, 0)  # Gold
powerbi_yellow = RGBColor(242, 200, 17)  # Power BI yellow
dark_text = RGBColor(26, 26, 26)
white = RGBColor(255, 255, 255)

def add_title_slide():
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background gradient effect
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = green_primary
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "FATHALLAH ELMASRI"
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = white
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle with Power BI emphasis
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(14), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.add_paragraph()
    p.text = "Senior Power BI Developer & Analytics Expert"
    p.font.size = Pt(32)
    p.font.color.rgb = powerbi_yellow
    p.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(14), Inches(2))
    contact_frame = contact_box.text_frame
    contact_frame.text = "📧 fathallah.elmasri@gmail.com | 📱 +966 593834672\n"
    contact_frame.add_paragraph().text = "🔗 LinkedIn | GitHub | Tableau Portfolio"
    for paragraph in contact_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = white
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Add Power BI badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6), Inches(7.5), Inches(4), Inches(0.8))
    badge.fill.solid()
    badge.fill.fore_color.rgb = powerbi_yellow
    badge.text_frame.text = "⚡ 6+ Years Power BI Excellence"
    badge.text_frame.paragraphs[0].font.size = Pt(18)
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].font.color.rgb = dark_text
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_powerbi_expertise_slide():
    slide_layout = prs.slide_layouts[5]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Power BI Center of Excellence"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = green_primary
    
    # Create content grid
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Add Power BI specific skills
    skills = [
        "⚡ Advanced DAX & Complex Measures",
        "📊 Executive Dashboard Development",
        "🎨 Custom Visuals & Themes",
        "⚙️ Data Modeling & Star Schemas",
        "🔄 Power Query/M Transformations",
        "📱 Power BI Service & Mobile",
        "🚀 Performance Optimization",
        "🔐 Row-Level Security (RLS)",
        "📈 Paginated Reports (SSRS)",
        "🔗 API Integration & Embedding"
    ]
    
    for skill in skills:
        p = tf.add_paragraph()
        p.text = skill
        p.font.size = Pt(22)
        p.level = 0
        p.font.bold = True if "Power BI" in skill else False

def add_stats_slide():
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Impact & Achievements"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = green_primary
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Stats grid
    stats = [
        ("100+", "Power BI\nDashboards"),
        ("30%", "Decision Speed\nImprovement"),
        ("500K+", "Daily Users\nServed"),
        ("60%", "Report Load\nOptimization"),
        ("15+", "Enterprise\nProjects")
    ]
    
    x_positions = [1.5, 4.5, 7.5, 10.5, 13.5]
    
    for i, (value, label) in enumerate(stats):
        # Stat box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                     Inches(x_positions[i]), Inches(3), 
                                     Inches(2), Inches(3))
        box.fill.solid()
        box.fill.fore_color.rgb = powerbi_yellow if i % 2 == 0 else RGBColor(240, 240, 240)
        
        # Add text
        tf = box.text_frame
        tf.clear()
        p1 = tf.add_paragraph()
        p1.text = value
        p1.font.size = Pt(36)
        p1.font.bold = True
        p1.font.color.rgb = dark_text if i % 2 == 0 else green_primary
        p1.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(16)
        p2.font.color.rgb = dark_text
        p2.alignment = PP_ALIGN.CENTER

def add_project_slide(title_text, company, tech_stack, achievements):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = green_primary
    
    # Company subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = company
    subtitle_frame.paragraphs[0].font.size = Pt(20)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
    subtitle_frame.paragraphs[0].font.italic = True
    
    # Tech stack badges
    tech_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(14), Inches(0.8))
    tech_frame = tech_box.text_frame
    tech_text = " | ".join(tech_stack)
    tech_frame.text = f"🔧 {tech_text}"
    tech_frame.paragraphs[0].font.size = Pt(16)
    tech_frame.paragraphs[0].font.color.rgb = powerbi_yellow if "Power BI" in tech_text else dark_text
    tech_frame.paragraphs[0].font.bold = True
    
    # Achievements
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    for achievement in achievements:
        p = tf.add_paragraph()
        p.text = f"✓ {achievement}"
        p.font.size = Pt(20)
        p.level = 0

def add_vision_slide():
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = green_primary
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Vision for Saudi Ministry of Industry"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = powerbi_yellow
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Vision statement
    vision_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(1.5))
    vision_frame = vision_box.text_frame
    vision_frame.text = "Establishing a World-Class Power BI Center of Excellence\nfor Industrial Analytics & Decision Intelligence"
    for paragraph in vision_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = white
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Three pillars
    pillars = [
        ("Foundation Phase", ["Power BI Governance", "Executive Dashboards", "Quick Wins", "Training Program"]),
        ("Expansion Phase", ["Self-Service Analytics", "Cross-Sector Platform", "Mobile Analytics", "Advanced DAX"]),
        ("Vision 2030", ["Industrial KPIs", "Sustainability Metrics", "Export Analytics", "Innovation Index"])
    ]
    
    x_positions = [1.5, 6, 10.5]
    
    for i, (pillar_title, items) in enumerate(pillars):
        # Pillar box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(x_positions[i]), Inches(4),
                                     Inches(4), Inches(4))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255) if i == 1 else RGBColor(240, 240, 240)
        box.fill.transparency = 0.1
        
        tf = box.text_frame
        tf.clear()
        
        # Pillar title
        p_title = tf.add_paragraph()
        p_title.text = pillar_title
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = powerbi_yellow if i == 1 else green_primary
        p_title.alignment = PP_ALIGN.CENTER
        
        # Items
        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = dark_text

def add_closing_slide():
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = green_primary
    
    # Thank you message
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "Thank You"
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = white
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Message
    msg_box = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(12), Inches(2))
    msg_frame = msg_box.text_frame
    msg_frame.text = "Ready to Transform Industrial Analytics with Power BI Excellence"
    msg_frame.paragraphs[0].font.size = Pt(24)
    msg_frame.paragraphs[0].font.color.rgb = powerbi_yellow
    msg_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Contact
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(14), Inches(1.5))
    contact_frame = contact_box.text_frame
    contact_frame.text = "Fathallah Elmasri\n📧 fathallah.elmasri@gmail.com | 📱 +966 593834672"
    for paragraph in contact_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = white
        paragraph.alignment = PP_ALIGN.CENTER

# Create all slides
add_title_slide()
add_powerbi_expertise_slide()
add_stats_slide()

# Add project slides
add_project_slide(
    "Enterprise Real Estate Analytics Platform",
    "AVIV Group | Data Model Manager | 2023-Present",
    ["Power BI Premium", "Advanced DAX", "Composite Models", "Snowflake"],
    [
        "15+ Executive Power BI dashboards with real-time KPIs",
        "Row-level security for multi-tenant data access",
        "60% report load time optimization through DAX tuning",
        "Custom Power Query functions for complex transformations",
        "Automated refresh pipelines using Power BI REST APIs"
    ]
)

add_project_slide(
    "Global Sales Intelligence Dashboard",
    "SHIJI Group | Senior BI Analyst | 2022-2023",
    ["Power BI Pro", "DAX Studio", "Paginated Reports", "Salesforce"],
    [
        "Sales funnel with advanced DAX conversion analysis",
        "Dynamic forecasting models with 85% accuracy",
        "Incremental refresh for 10M+ records",
        "Mobile-optimized C-suite dashboards",
        "Salesforce integration via Power BI dataflows"
    ]
)

add_project_slide(
    "Solar Energy Financial Analytics",
    "Enpal | Business Intelligence Analyst | 2023",
    ["Power BI Embedded", "Custom Visuals", "Power Apps", "Azure SQL"],
    [
        "Power BI embedded analytics for customer portal",
        "Complex DAX for ROI and payback calculations",
        "Custom visuals for solar efficiency metrics",
        "Real-time streaming datasets",
        "Power Apps integration for field analytics"
    ]
)

add_vision_slide()
add_closing_slide()

# Save presentation
output_path = "PowerBI_Portfolio_KSA_Ministry.pptx"
prs.save(output_path)
print(f"✅ PowerPoint presentation saved as: {output_path}")
