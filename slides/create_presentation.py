"""
PowerPoint Portfolio Generator for Senior Power BI Developer
Saudi Ministry of Industry Application
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_portfolio_presentation():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Define Ministry colors
    ministry_green = RGBColor(46, 125, 50)  # #2E7D32
    ministry_dark = RGBColor(27, 94, 32)    # #1B5E20
    ministry_gold = RGBColor(201, 169, 97)  # #C9A961
    accent_blue = RGBColor(25, 118, 210)    # #1976D2
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Senior Power BI Developer Portfolio"
    subtitle.text = "Fathallah Elmasri\nBusiness Intelligence Expert\nPrepared for: Saudi Ministry of Industry\nJanuary 2025"
    
    # Slide 2: Table of Contents
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Portfolio Contents"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. Executive Summary & Achievements\n"
    tf.text += "2. Core Competencies & Technical Stack\n"
    tf.text += "3. Key Project: Real Estate Intelligence Platform\n"
    tf.text += "4. Key Project: Financial Performance Suite\n"
    tf.text += "5. Key Project: Sales Intelligence Platform\n"
    tf.text += "6. Technical Demonstrations\n"
    tf.text += "7. Professional Timeline\n"
    tf.text += "8. Vision for Ministry of Industry\n"
    tf.text += "9. Contact & Digital Portfolio"
    
    # Slide 3: Executive Summary
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Executive Summary"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Senior Power BI Developer | 6+ Years Experience\n\n"
    
    p = tf.add_paragraph()
    p.text = "Key Achievements:"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 30% improvement in sales visibility (SHIJI Group)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 20% productivity increase through BI optimization"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 500+ users served across enterprise platforms"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• €2B+ portfolio monitored in real-time"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nMultilingual: Arabic (Native) | English (Fluent) | German (B1)"
    p.level = 0
    
    # Slide 4: Core Competencies
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Core Competencies"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Power BI Ecosystem Mastery\n"
    
    p = tf.add_paragraph()
    p.text = "• Desktop, Service, Paginated Reports, Power Query, DAX"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nData Architecture & Modeling"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Star Schema, Snowflake, Azure, dbt, ETL/ELT"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nAnalytics & Visualization"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• KPIs, Predictive Models, Mobile Optimization"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nSupporting Technologies"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• SQL, Python, Tableau, Excel (Advanced), Git"
    p.level = 1
    
    # Slide 5: AVIV Group Project
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Real Estate Intelligence Platform"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "AVIV Group | Data Model Manager | 2023-Present\n\n"
    
    p = tf.add_paragraph()
    p.text = "Challenge: Unified BI across multiple real estate platforms"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "\nSolution:"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Enterprise Power BI suite serving 500+ users"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 10M+ transactions processed monthly"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 150+ DAX measures for complex calculations"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nImpact:"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 40% faster insights | €2B+ portfolio | Audit-ready"
    p.level = 1
    
    # Slide 6: Enpal Project
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Financial Performance Dashboard Suite"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Enpal | BI Analyst | June-Sept 2023\n\n"
    
    p = tf.add_paragraph()
    p.text = "Challenge: Real-time financial tracking for solar energy growth"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "\nSolution:"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Executive dashboard tracking €100M+ revenue"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Complex DAX for CAC, LTV, payback calculations"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 95% forecast accuracy predictive models"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nImpact:"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Real-time visibility | 25% better forecasting | 20 hrs/week saved"
    p.level = 1
    
    # Slide 7: SHIJI Project
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Sales Intelligence & Strategy Platform"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "SHIJI Group | BI Analyst | Sept 2022-April 2023\n\n"
    
    p = tf.add_paragraph()
    p.text = "Challenge: Optimize global sales operations"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "\nSolution:"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 30% improvement in sales visibility"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 24-month projection models"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Territory analysis for 50+ countries"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nImpact:"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• €5M+ opportunities | 15% shorter sales cycle"
    p.level = 1
    
    # Slide 8: Technical Demonstrations
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Technical Demonstrations"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Customer Journey Analytics (Enpal Case)\n"
    
    p = tf.add_paragraph()
    p.text = "• Multi-stage funnel visualization"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Lead response optimization"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Predictive installation models"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nRevenue Analytics Platform (SHIJI Case)"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Time-series decomposition"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Customer lifetime value"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Churn prediction (85% accuracy)"
    p.level = 1
    
    # Slide 9: Professional Timeline
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Professional Timeline"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "2023-Present: Data Model Manager - AVIV Group\n"
    tf.text += "Leading enterprise BI strategy for real estate platform\n\n"
    
    tf.text += "2023: BI Analyst - Enpal\n"
    tf.text += "Financial dashboards for solar energy sector\n\n"
    
    tf.text += "2022-2023: BI Analyst - SHIJI Group\n"
    tf.text += "Sales intelligence and strategy platform\n\n"
    
    tf.text += "2019-2022: Team Manager - Sykes Enterprises\n"
    tf.text += "BI and operations reporting for 12-18 member team\n\n"
    
    tf.text += "Education:\n"
    tf.text += "• M.Sc. Geodesy & Geoinformation (2025) - TU Berlin\n"
    tf.text += "• B.Sc. Civil Engineering (2012) - Mansoura University"
    
    # Slide 10: Vision for Ministry
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Vision for Saudi Ministry of Industry"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Supporting Vision 2030 through Data Excellence\n\n"
    
    p = tf.add_paragraph()
    p.text = "30-Day Plan:"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Assess current BI landscape"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Establish Power BI governance"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Design prototype dashboards"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n90-Day Plan:"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Deploy production dashboards"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Implement self-service analytics"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Train ministry staff"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n1-Year Vision:"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Comprehensive industrial intelligence platform"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Predictive analytics for sector planning"
    p.level = 1
    
    # Slide 11: Contact & Portfolio
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Contact & Digital Portfolio"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Fathallah Elmasri\n\n"
    
    tf.text += "📧 fathallah.elmasri@gmail.com\n"
    tf.text += "📱 (+966) 593834672 | (+49) 15510197596\n\n"
    
    tf.text += "Digital Portfolio:\n"
    tf.text += "🔗 Tableau: public.tableau.com/app/profile/fathallah.elmasri\n"
    tf.text += "🔗 GitHub: github.com/Elmasri-Fathallah\n"
    tf.text += "🔗 LinkedIn: linkedin.com/in/fathallah-elmasri\n\n"
    
    tf.text += "Thank you for considering my application.\n"
    tf.text += "I look forward to contributing to the Kingdom's\n"
    tf.text += "industrial transformation through data-driven excellence."
    
    # Save presentation
    output_path = r"C:\Users\fatah\OneDrive\Desktop\MY Files\Job_Application\Interviews-Technical_Tasks\KSAPortfolia\portfolio_project\slides\Portfolio_Deck.pptx"
    prs.save(output_path)
    print(f"PowerPoint presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    try:
        output = create_portfolio_presentation()
        print(f"Success! Presentation created at: {output}")
    except Exception as e:
        print(f"Error creating presentation: {e}")
        print("\nPlease ensure python-pptx is installed:")
        print("pip install python-pptx")
